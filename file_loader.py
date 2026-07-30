"""
Multi-modal file parsing module.
Extracts text from common document formats for profile ingestion.
"""
import os
import zipfile
import xml.etree.ElementTree as ET
import logging
from typing import Tuple
logger = logging.getLogger(__name__)
class FileLoader:
    """Document parser supporting text, PPTX, DOCX, PDF, XLSX."""
    TEXT_EXTS = ('.txt', '.md', '.json', '.csv', '.log', '.py', '.js', '.ts', '.html', '.css')
    def extract(self, file_path: str) -> Tuple[str, str]:
        """Extract file content, returns (text_content, format_label)."""
        ext = os.path.splitext(file_path)[1].lower()
        if ext in self.TEXT_EXTS:
            return self._read_text(file_path), f"text:{ext}"
        if ext == '.pptx':
            return self._read_pptx(file_path), "pptx:slides"
        if ext == '.docx':
            return self._read_docx(file_path), "docx:paragraphs"
        if ext == '.pdf':
            return self._read_pdf(file_path), "pdf:pages"
        if ext in ('.xlsx', '.xlsm'):
            return self._read_xlsx(file_path), "xlsx:sheets"
        raise ValueError(f"Unsupported format: {ext}")
    def _read_text(self, path: str) -> str:
        try:
            with open(path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            with open(path, 'r', encoding='gb18030', errors='replace') as f:
                return f.read()
    def _read_pptx(self, path: str) -> str:
        try:
            from pptx import Presentation
        except ImportError:
            return self._read_pptx_zip_fallback(path, "python-pptx not installed")
        prs = Presentation(path)
        lines = [f"# PPTX Parse Result: {os.path.basename(path)}", f"Total slides: {len(prs.slides)}"]
        for idx, slide in enumerate(prs.slides, 1):
            lines.append(f"\n## Slide {idx}")
            texts = []
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False) and shape.text_frame:
                    t = shape.text_frame.text.strip()
                    if t: texts.append(t)
            lines.extend(texts)
        return "\n".join(lines)
    def _read_pptx_zip_fallback(self, path: str, reason: str) -> str:
        lines = [f"# PPTX ZIP fallback: {os.path.basename(path)}", reason]
        try:
            ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
            with zipfile.ZipFile(path) as z:
                slides = sorted(n for n in z.namelist() if n.startswith('ppt/slides/slide') and n.endswith('.xml'))
                lines.append(f"Total slides: {len(slides)}")
                for idx, name in enumerate(slides, 1):
                    root = ET.fromstring(z.read(name))
                    texts = [t.text for t in root.findall('.//a:t', ns) if t.text]
                    lines.append(f"\n## Slide {idx}\n" + "\n".join(texts))
        except Exception as e:
            lines.append(f"PPTX fallback failed: {e}")
        return "\n".join(lines)
    def _read_docx(self, path: str) -> str:
        try:
            from docx import Document
            doc = Document(path)
            lines = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n".join(lines)
        except ImportError:
            return "DOCX parse failed: python-docx not installed"
        except Exception as e:
            return f"DOCX parse failed: {e}"
    def _read_pdf(self, path: str) -> str:
        try:
            import pypdf
            reader = pypdf.PdfReader(path)
            lines = [f"# PDF Parse Result: {os.path.basename(path)}", f"Total pages: {len(reader.pages)}"]
            for i, page in enumerate(reader.pages, 1):
                lines.append(f"\n## Page {i}\n{page.extract_text() or ''}")
            return "\n".join(lines)
        except ImportError:
            return "PDF parse failed: pypdf not installed"
        except Exception as e:
            return f"PDF parse failed: {e}"
    def _read_xlsx(self, path: str) -> str:
        try:
            import openpyxl
            wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
            lines = [f"# Excel Parse Result: {os.path.basename(path)}"]
            for ws in wb.worksheets:
                lines.append(f"\n## Sheet: {ws.title}")
                for row in ws.iter_rows(values_only=True):
                    vals = ["" if v is None else str(v) for v in row]
                    if any(v.strip() for v in vals):
                        lines.append(" | ".join(vals))
            return "\n".join(lines)
        except ImportError:
            return "Excel parse failed: openpyxl not installed"
        except Exception as e:
            return f"Excel parse failed: {e}"
